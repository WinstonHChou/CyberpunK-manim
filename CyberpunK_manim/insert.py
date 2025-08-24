from pptx import Presentation

def copy_slide(source_slide, target_prs):
    slide_layout = target_prs.slide_layouts[0]  # Use blank layout
    new_slide = target_prs.slides.add_slide(slide_layout)
    for shape in source_slide.shapes:
        if shape.shape_type == 1:  # Placeholder
            continue
        el = shape.element
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    return new_slide

def extract_slides(pptx_path):
    prs = Presentation(pptx_path)
    return [slide for slide in prs.slides]

def arrange_and_insert_slides(source_slides, target_pptx_path, order, output_path):
    target_prs = Presentation(target_pptx_path)
    for idx in order:
        copy_slide(source_slides[idx], target_prs)
    target_prs.save(output_path)

if __name__ == "__main__":
    # Example usage
    source_pptx = "presentation2.pptx"
    target_pptx = "presentation1.pptx"
    output_pptx = "combined_custom_order.pptx"

    # Extract slides from source presentation
    slides = extract_slides(source_pptx)

    # Developer chooses the order, e.g., [2, 0, 1] means 3rd, 1st, then 2nd slide
    chosen_order = [2, 0, 1]  # Change as needed

    # Insert selected slides into target presentation in chosen order
    arrange_and_insert_slides(slides, target_pptx, chosen_order, output_pptx)