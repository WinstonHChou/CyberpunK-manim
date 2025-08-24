from pptx import Presentation
import os

# Parent directory of the current script
dir_path = os.path.dirname(os.path.abspath(__file__))

prs = Presentation(f"{dir_path}/CyberpunK-theme.pptx")

cover = prs.slides[0]
cover.shapes.title.text = 'Cover Slide Title'
cover.shapes.placeholders[1].text = 'Cover Slide Content'

bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')