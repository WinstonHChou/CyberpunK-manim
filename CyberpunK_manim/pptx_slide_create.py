from CyberpunK_manim import Color, render_and_extract, add_animation, get_pptx_presentation
# requires `pip install python-pptx`
from pptx.util import *
from pptx.enum.text import PP_ALIGN

def edit_cover(prs, title, subtitle):
    ## Cover Slide
    cover = prs.slides[0]
    cover.shapes.title.text = title
    cover_subtitle = cover.shapes.placeholders[1].text_frame.paragraphs[0].add_run()
    cover_subtitle.text = subtitle
    return cover

def edit_cover_with_subtitle_color(prs, title, subtitle, color=Color.CK_LIGHT_BLUE, bold=True):
    ## Cover Slide
    cover = prs.slides[0]
    cover.shapes.title.text = title
    cover_subtitle = cover.shapes.placeholders[1].text_frame.paragraphs[0].add_run()
    cover_subtitle.text = subtitle
    cover_subtitle.font.bold = bold
    cover_subtitle.font.color.rgb = color
    return cover

def add_title_slide(prs, title_slide_layout, title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle

    return slide


def add_bullet_slide(prs, title, content=None, picture_path=None, picture_label=None):
    bullet_slide_layout = prs.slide_layouts[1]  # Use the "Title and Content" layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title

    if content:
        tf = body_shape.text_frame
        tf.text = content
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left alignment

        # title_shape = shapes.title
        # body_shape = shapes.placeholders[1]
        # title_shape.text = 'What is a vector'
        # tf = body_shape.text_frame
        # tf.text = 'Physics Perspective'
        # p = tf.add_paragraph()
        # p.text = 'Computer Science Perspective'
        # p.level = 0
        # p = tf.add_paragraph()
        # p.text = 'Mathematician’s Perspective'
        # p.level = 0

    if picture_path:
        # Add the image
        img_left = Inches(5.75)
        img_top = Inches(1.5)
        img_width = prs.slide_width * 9 / 16
        img_height = prs.slide_height * 9 / 16
        pic = shapes.add_picture(picture_path, img_left, img_top, img_width, img_height)

        if picture_label:
            # Add a left-aligned textbox label below the image
            label_left = img_left
            label_top = img_top + img_height
            label_width = img_width
            label_height = Inches(0.6)
            textBox = shapes.add_textbox(label_left, label_top, label_width, label_height)
            tf = textBox.text_frame
            tf.text = picture_label
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left alignment

    return slide, title_shape, body_shape